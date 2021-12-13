from pptx import Presentation
from paragraphs_extractor.file_iterator_interface import FileIteratorInterface

class PPTXIterator(FileIteratorInterface):
    def __init__(self, filename):
        self.filename = filename
        self.paragraphs = []
        
        prs = Presentation(filename)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    cleaned_text = shape.text.replace('\n', '')
                    if cleaned_text:
                        self.paragraphs.append(cleaned_text)
    